from __future__ import annotations

import base64
import shutil
import uuid
import zipfile
from pathlib import Path
from typing import List

import fitz  # PyMuPDF
from bs4 import BeautifulSoup
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from markdownify import markdownify as mdify
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parent
OUTPUTS_DIR = BASE_DIR / "outputs"
OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)

app = FastAPI(title="pptToMD")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"] ,
    allow_headers=["*"],
)


def _safe_filename(name: str) -> str:
    return "".join(c for c in name if c.isalnum() or c in ("-", "_", "."))


def _write_image_bytes(images_dir: Path, filename: str, data: bytes) -> str:
    images_dir.mkdir(parents=True, exist_ok=True)
    file_path = images_dir / filename
    file_path.write_bytes(data)
    return f"images/{filename}"


def convert_pptx(pptx_path: Path, images_dir: Path) -> str:
    prs = Presentation(str(pptx_path))
    md_lines: List[str] = []
    img_index = 1

    for slide_idx, slide in enumerate(prs.slides, start=1):
        md_lines.append(f"# Slide {slide_idx}")

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                filename = f"slide{slide_idx}_img{img_index}.{ext}"
                rel_path = _write_image_bytes(images_dir, filename, image.blob)
                md_lines.append(f"![image]({rel_path})")
                img_index += 1
                continue

            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    md_lines.append(text)
                continue

            if getattr(shape, "has_table", False):
                table = shape.table
                rows = table.rows
                cols = table.columns
                if rows and cols:
                    header_cells = [table.cell(0, c).text.strip() for c in range(len(cols))]
                    md_lines.append("| " + " | ".join(header_cells) + " |")
                    md_lines.append("| " + " | ".join(["---"] * len(cols)) + " |")
                    for r in range(1, len(rows)):
                        row_cells = [table.cell(r, c).text.strip() for c in range(len(cols))]
                        md_lines.append("| " + " | ".join(row_cells) + " |")

        md_lines.append("")

    return "\n".join(md_lines).strip() + "\n"


def convert_pdf(pdf_path: Path, images_dir: Path) -> str:
    doc = fitz.open(str(pdf_path))
    md_lines: List[str] = []
    img_index = 1

    for page_index in range(len(doc)):
        page = doc[page_index]
        md_lines.append(f"# Page {page_index + 1}")

        text = page.get_text().strip()
        if text:
            md_lines.append(text)

        for img in page.get_images(full=True):
            xref = img[0]
            extracted = doc.extract_image(xref)
            ext = extracted.get("ext", "png")
            image_bytes = extracted.get("image", b"")
            if not image_bytes:
                continue
            filename = f"page{page_index + 1}_img{img_index}.{ext}"
            rel_path = _write_image_bytes(images_dir, filename, image_bytes)
            md_lines.append(f"![image]({rel_path})")
            img_index += 1

        md_lines.append("")

    return "\n".join(md_lines).strip() + "\n"


def _handle_html_image(img_tag, html_dir: Path, images_dir: Path, index: int) -> None:
    src = img_tag.get("src", "").strip()
    if not src:
        return

    if src.startswith("data:image/"):
        header, b64data = src.split(",", 1)
        mime = header.split(";")[0]
        ext = mime.split("/")[-1] or "png"
        filename = f"image{index}.{ext}"
        rel_path = _write_image_bytes(images_dir, filename, base64.b64decode(b64data))
        img_tag["src"] = rel_path
        return

    if src.startswith("http://") or src.startswith("https://"):
        return

    src_path = (html_dir / src).resolve()
    if src_path.exists() and src_path.is_file():
        ext = src_path.suffix.lstrip(".") or "png"
        filename = f"image{index}.{ext}"
        images_dir.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(src_path, images_dir / filename)
        img_tag["src"] = f"images/{filename}"


def convert_html(html_path: Path, images_dir: Path) -> str:
    html_dir = html_path.parent
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8", errors="ignore"), "html.parser")

    img_index = 1
    for img in soup.find_all("img"):
        _handle_html_image(img, html_dir, images_dir, img_index)
        img_index += 1

    return mdify(str(soup), heading_style="ATX").strip() + "\n"


@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/convert")
async def convert(file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")

    filename = _safe_filename(file.filename)
    ext = Path(filename).suffix.lower()

    if ext not in {".pptx", ".pdf", ".html", ".htm"}:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    job_id = uuid.uuid4().hex
    job_dir = OUTPUTS_DIR / job_id
    images_dir = job_dir / "images"
    job_dir.mkdir(parents=True, exist_ok=True)

    input_path = job_dir / f"input{ext}"
    with input_path.open("wb") as f:
        content = await file.read()
        f.write(content)

    if ext == ".pptx":
        markdown = convert_pptx(input_path, images_dir)
    elif ext == ".pdf":
        markdown = convert_pdf(input_path, images_dir)
    else:
        markdown = convert_html(input_path, images_dir)

    markdown_path = job_dir / "output.md"
    markdown_path.write_text(markdown, encoding="utf-8")

    images = []
    if images_dir.exists():
        images = sorted([f"images/{p.name}" for p in images_dir.iterdir() if p.is_file()])

    return {
        "job_id": job_id,
        "markdown": markdown,
        "markdown_path": str(markdown_path),
        "images": images,
    }


@app.get("/download/{job_id}/markdown")
def download_markdown(job_id: str):
    job_dir = OUTPUTS_DIR / job_id
    markdown_path = job_dir / "output.md"
    if not markdown_path.exists():
        raise HTTPException(status_code=404, detail="Not found")
    return FileResponse(markdown_path, filename="output.md")


@app.get("/download/{job_id}/bundle")
def download_bundle(job_id: str):
    job_dir = OUTPUTS_DIR / job_id
    markdown_path = job_dir / "output.md"
    if not markdown_path.exists():
        raise HTTPException(status_code=404, detail="Not found")

    zip_path = job_dir / "output_bundle.zip"
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as zip_file:
        zip_file.write(markdown_path, arcname="output.md")

        images_dir = job_dir / "images"
        if images_dir.exists():
            for image_path in sorted(images_dir.iterdir()):
                if image_path.is_file():
                    zip_file.write(image_path, arcname=f"images/{image_path.name}")

    return FileResponse(zip_path, filename="output_bundle.zip")


@app.get("/download/{job_id}/assets/{filename}")
def download_asset(job_id: str, filename: str):
    safe_name = _safe_filename(filename)
    job_dir = OUTPUTS_DIR / job_id
    asset_path = job_dir / "images" / safe_name
    if not asset_path.exists():
        raise HTTPException(status_code=404, detail="Not found")
    return FileResponse(asset_path, filename=safe_name)
